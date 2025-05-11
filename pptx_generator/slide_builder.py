#!/usr/bin/env python3
import os
from pptx import Presentation
from pptx.util import Pt
from reco.models import DeckData

def build_ppt(deck: DeckData, output_path: str):
    # Vérifie la présence du template, fallback à une présentation vide
    template_path = "pptx_generator/templates/base.pptx"
    if os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # 1. Brief Reminder
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. Brief Reminder"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for obj in deck.brief_reminder.objectives:
        p = tf.add_paragraph()
        p.text = f"• {obj}"
        p.level = 0
    p = tf.add_paragraph()
    p.text = deck.brief_reminder.internal_reformulation
    p.level = 0

    # 2. Brand Overview – description
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2. Brand Overview"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for paragraph in deck.brand_overview.description_paragraphs:
        p = tf.add_paragraph()
        p.text = paragraph
        p.level = 0

    # 2b. Competitive Positioning & Persona (à compléter)
    # TODO: ajouter graphique/tableau pour deck.brand_overview.competitive_positioning
    # TODO: ajouter slide pour persona (deck.brand_overview.persona)

    # 3. State of Play
    for idx, section in enumerate(deck.state_of_play, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"3.{idx} State of Play – {section.theme}"
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for ev in section.evidence:
            p = tf.add_paragraph()
            p.text = f"• {ev.title}: {ev.desc}"
            p.level = 0

    # 4. Ideas
    for idx, idea in enumerate(deck.ideas, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"4. Idea #{idx}: {idea.title}"
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = idea.body
        p.level = 0

    # 5. Timeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "5. Timeline"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for m in deck.timeline:
        p = tf.add_paragraph()
        p.text = f"{m.date}: {m.milestone}"
        p.level = 0

    # 6. Budget
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "6. Budget"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for b in deck.budget:
        p = tf.add_paragraph()
        p.text = f"{b.item}: €{b.amount}"
        p.level = 0

    # Sauvegarde du PPT
    prs.save(output_path)
